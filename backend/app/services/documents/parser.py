"""Extract plain text from uploaded documents."""
import io
import logging

log = logging.getLogger(__name__)

_DOCX = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
_XLSX = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
_PPTX = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

SUPPORTED_MIME_TYPES = {
    "application/pdf",
    _DOCX,
    _XLSX,
    _PPTX,
    "text/plain",
    "text/markdown",
}


def parse(content: bytes, mime_type: str) -> str:
    if mime_type == "application/pdf":
        return _parse_pdf(content)
    if mime_type == _DOCX:
        return _parse_docx(content)
    if mime_type == _XLSX:
        return _parse_xlsx(content)
    if mime_type == _PPTX:
        return _parse_pptx(content)
    if mime_type in ("text/plain", "text/markdown"):
        return _parse_txt(content)
    raise ValueError(f"Unsupported MIME type: {mime_type}")


def _parse_pdf(content: bytes) -> str:
    import fitz  # PyMuPDF

    doc = fitz.open(stream=content, filetype="pdf")
    pages = []
    for page in doc:
        pages.append(page.get_text())
    doc.close()
    return "\n".join(pages).strip()


def _parse_docx(content: bytes) -> str:
    from docx import Document

    doc = Document(io.BytesIO(content))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n".join(paragraphs).strip()


def _parse_xlsx(content: bytes) -> str:
    import openpyxl

    wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    parts: list[str] = []
    for ws in wb.worksheets:
        parts.append(f"# Sheet: {ws.title}")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                parts.append(" | ".join(cells))
    wb.close()
    return "\n".join(parts).strip()


def _parse_pptx(content: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(content))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts.append(f"# Slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text.strip())
    return "\n".join(parts).strip()


def _parse_txt(content: bytes) -> str:
    for encoding in ("utf-8", "utf-8-sig", "latin-1"):
        try:
            return content.decode(encoding).strip()
        except UnicodeDecodeError:
            continue
    return content.decode("latin-1", errors="replace").strip()
